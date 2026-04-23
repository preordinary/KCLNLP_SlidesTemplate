"""Build the Tree Yellow Texture demo deck.

Run: python templates/tree_yellow_texture/build.py

Convention: slide titles live inside the cream header band (left),
freeing the body for content. Use the ``title=`` argument on
``slide_content`` / ``slide_two_column`` — do not add a second title
inside the body.
"""
from __future__ import annotations

import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))

import theme as T

OUTPUT = Path(__file__).resolve().parent / "KCLNLP_TreeYellow.pptx"


def rgb(tup):
    return RGBColor(*tup)


CROP_DIR = Path(__file__).resolve().parents[2] / "assets" / "logos" / "_trimmed"


def _trimmed_logo(path: Path) -> Path:
    """Return a cached copy of the logo cropped to its visible (alpha) bbox.

    Raw PNGs carry different amounts of transparent padding, so scaling them
    to the same image height produces different visible heights. Trimming to
    the alpha bbox first normalises visible height and inter-logo gaps.
    """
    CROP_DIR.mkdir(parents=True, exist_ok=True)
    dst = CROP_DIR / path.name
    if dst.exists() and dst.stat().st_mtime >= path.stat().st_mtime:
        return dst
    with Image.open(path) as im:
        im = im.convert("RGBA")
        bbox = im.getbbox()
        if bbox:
            im = im.crop(bbox)
        im.save(dst, "PNG", optimize=True)
    return dst


def _logo_width(path, height_emu):
    with Image.open(path) as im:
        w, h = im.size
    return int(height_emu * (w / h))


def add_header(slide, title: str | None = None):
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, T.SLIDE_W, T.HEADER_H
    )
    band.fill.solid()
    band.fill.fore_color.rgb = rgb(T.CREAM)
    band.line.fill.background()

    trimmed = [_trimmed_logo(p) for p in T.LOGOS]
    top = (T.HEADER_H - T.LOGO_H) // 2
    right = T.SLIDE_W - T.EDGE_PAD
    widths = [_logo_width(p, T.LOGO_H) for p in trimmed]
    total = sum(widths) + T.LOGO_GAP * (len(trimmed) - 1)
    x = right - total
    for i, path in enumerate(trimmed):
        slide.shapes.add_picture(str(path), x, top, height=T.LOGO_H)
        x += widths[i] + (T.LOGO_GAP if i < len(trimmed) - 1 else 0)

    if title:
        logos_left = right - total
        title_w = logos_left - T.EDGE_PAD - Inches(0.3)
        box = slide.shapes.add_textbox(
            T.EDGE_PAD, 0, title_w, T.HEADER_H
        )
        tf = box.text_frame
        tf.margin_left = Emu(0); tf.margin_right = Emu(0)
        tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT

        now_run = p.add_run()
        now_run.text = "Now:  "
        now_run.font.name = T.FONT_HEAD
        now_run.font.size = T.SZ_HEADER_NOW
        now_run.font.bold = True
        now_run.font.italic = True
        now_run.font.color.rgb = rgb(T.ORANGE)

        title_run = p.add_run()
        title_run.text = title
        title_run.font.name = T.FONT_HEAD
        title_run.font.size = T.SZ_HEADER_TITLE
        title_run.font.bold = True
        title_run.font.italic = True
        title_run.font.color.rgb = rgb(T.INK)

    rule_y = T.HEADER_H
    rule = slide.shapes.add_connector(
        1, 0, rule_y, T.SLIDE_W, rule_y
    )
    rule.line.color.rgb = rgb(T.ORANGE)
    rule.line.width = Emu(19050)


def add_title_accent(slide, x, y, h):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.12), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(T.ORANGE)
    bar.line.fill.background()


def _add_text(slide, left, top, width, height, text, *, font, size, color,
              bold=False, italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = size
    run.font.color.rgb = rgb(color)
    run.font.bold = bold
    run.font.italic = italic
    return box


def _add_bullets(slide, left, top, width, height, items, *, size=T.SZ_BODY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.name = T.FONT_BODY
        run.font.size = size
        run.font.color.rgb = rgb(T.INK)
    return box


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide)

    left = T.EDGE_PAD + Inches(0.2)
    title_top = Inches(2.7)

    add_title_accent(slide, left, title_top + Inches(0.1), Inches(1.4))
    _add_text(slide, left + Inches(0.35), title_top,
              Inches(11), Inches(1.6),
              "Talk Title Goes Here",
              font=T.FONT_HEAD, size=T.SZ_TITLE, color=T.INK, bold=True)

    _add_text(slide, left + Inches(0.35), title_top + Inches(1.5),
              Inches(11), Inches(0.6),
              "A concise subtitle or one-line abstract",
              font=T.FONT_BODY, size=T.SZ_SUBTITLE, color=T.MUTED, italic=True)

    rule_y = title_top + Inches(2.3)
    rule = slide.shapes.add_connector(
        1, left + Inches(0.35), rule_y,
        left + Inches(2.2), rule_y
    )
    rule.line.color.rgb = rgb(T.ORANGE)
    rule.line.width = Emu(19050)

    _add_text(slide, left + Inches(0.35), rule_y + Inches(0.15),
              Inches(11), Inches(0.4),
              "Speaker Name  ·  Affiliation  ·  Month Year",
              font=T.FONT_BODY, size=Pt(16), color=T.INK)


def slide_title_centered(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide)

    body_top = T.HEADER_H
    body_h = T.SLIDE_H - body_top
    block_h = Inches(3.6)
    block_top = body_top + (body_h - block_h) // 2

    accent_w = Inches(1.6)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        (T.SLIDE_W - accent_w) // 2, block_top,
        accent_w, Inches(0.06),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(T.ORANGE)
    accent.line.fill.background()

    _add_text(slide, Inches(0), block_top + Inches(0.35),
              T.SLIDE_W, Inches(1.4),
              "Talk Title Goes Here",
              font=T.FONT_HEAD, size=T.SZ_TITLE, color=T.INK, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _add_text(slide, Inches(0), block_top + Inches(1.85),
              T.SLIDE_W, Inches(0.7),
              "A concise subtitle or one-line abstract",
              font=T.FONT_BODY, size=T.SZ_SUBTITLE, color=T.MUTED, italic=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    rule_y = block_top + Inches(2.75)
    rule_half = Inches(1.1)
    cx = T.SLIDE_W // 2
    rule = slide.shapes.add_connector(
        1, cx - rule_half, rule_y, cx + rule_half, rule_y
    )
    rule.line.color.rgb = rgb(T.ORANGE)
    rule.line.width = Emu(19050)

    _add_text(slide, Inches(0), rule_y + Inches(0.2),
              T.SLIDE_W, Inches(0.5),
              "Speaker Name  ·  Affiliation  ·  Month Year",
              font=T.FONT_BODY, size=Pt(16), color=T.INK,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)


def slide_section(prs, number="01", title="Section Heading"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title=title)

    left = T.EDGE_PAD + Inches(0.5)
    row_top = Inches(3.15)
    row_h = Inches(1.4)

    _add_text(slide, left, row_top, Inches(1.8), row_h,
              number, font=T.FONT_HEAD, size=Pt(56), color=T.ORANGE, bold=True,
              anchor=MSO_ANCHOR.MIDDLE)

    rule_x = left + Inches(1.9)
    rule = slide.shapes.add_connector(
        1, rule_x, row_top + Inches(0.25),
        rule_x, row_top + row_h - Inches(0.25)
    )
    rule.line.color.rgb = rgb(T.ORANGE)
    rule.line.width = Emu(12700)

    _add_text(slide, rule_x + Inches(0.3), row_top,
              Inches(8), row_h,
              title, font=T.FONT_HEAD, size=T.SZ_SECTION,
              color=T.INK, bold=True,
              anchor=MSO_ANCHOR.MIDDLE)

    _add_text(slide, rule_x + Inches(0.3), row_top + row_h + Inches(0.05),
              Inches(8), Inches(0.5),
              "An optional one-line summary of this section",
              font=T.FONT_BODY, size=Pt(18), color=T.MUTED, italic=True)


def slide_content(prs, title="Slide title", bullets=None):
    bullets = bullets or [
        "Lead with the main claim — what the reader should take away",
        "Supporting detail, ideally a number or concrete example",
        "A second supporting point that extends the first",
        "A fourth line of room — titles live in the header, so the body breathes",
        "Edge case, caveat, or the one thing not to forget",
    ]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title=title)

    body_left = T.EDGE_PAD + Inches(0.2)
    body_top = T.HEADER_H + Inches(0.35)
    body_w = T.SLIDE_W - body_left - T.EDGE_PAD
    body_h = T.SLIDE_H - body_top - Inches(0.55)

    _add_bullets(slide, body_left, body_top, body_w, body_h, bullets)

    _add_text(slide, body_left, T.SLIDE_H - Inches(0.45),
              body_w, Inches(0.35),
              "Footer / citation / page note",
              font=T.FONT_BODY, size=T.SZ_CAPTION, color=T.MUTED, italic=True)


def slide_two_column(prs, title="Slide title"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title=title)

    body_left = T.EDGE_PAD + Inches(0.2)
    body_top = T.HEADER_H + Inches(0.35)
    body_h = T.SLIDE_H - body_top - Inches(0.5)
    col_w = Inches(6.0)

    frame = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, body_left, body_top, col_w, body_h
    )
    frame.fill.background()
    frame.line.color.rgb = rgb(T.MUTED)
    frame.line.width = Emu(6350)
    _add_text(slide, body_left, body_top, col_w, body_h,
              "[ figure / diagram ]",
              font=T.FONT_BODY, size=Pt(18), color=T.MUTED, italic=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    text_x = body_left + col_w + Inches(0.5)
    text_w = T.SLIDE_W - text_x - T.EDGE_PAD
    _add_text(slide, text_x, body_top, text_w, Inches(0.5),
              "Caption or claim",
              font=T.FONT_HEAD, size=Pt(22), color=T.INK, bold=True)
    _add_bullets(slide, text_x, body_top + Inches(0.8),
                 text_w, body_h - Inches(0.8),
                 ["Explain what the figure shows",
                  "Call out the one feature worth noting",
                  "Relate it back to the slide's main point"],
                 size=Pt(18))


def slide_thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide)

    left = T.EDGE_PAD + Inches(0.2)
    y = Inches(3.0)

    add_title_accent(slide, left, y + Inches(0.15), Inches(1.6))
    _add_text(slide, left + Inches(0.35), y,
              Inches(10), Inches(1.6),
              "Thank you",
              font=T.FONT_HEAD, size=Pt(56), color=T.INK, bold=True)

    rule_y = y + Inches(1.75)
    rule = slide.shapes.add_connector(
        1, left + Inches(0.35), rule_y,
        left + Inches(2.2), rule_y
    )
    rule.line.color.rgb = rgb(T.ORANGE)
    rule.line.width = Emu(19050)

    _add_text(slide, left + Inches(0.35), rule_y + Inches(0.2),
              Inches(10), Inches(0.6),
              "Questions & discussion",
              font=T.FONT_BODY, size=Pt(22), color=T.MUTED, italic=True)

    _add_text(slide, left + Inches(0.35), rule_y + Inches(1.1),
              Inches(10), Inches(0.45),
              "speaker@kcl.ac.uk  ·  kclnlp.github.io",
              font=T.FONT_BODY, size=Pt(16), color=T.INK)


def slide_thanks_centered(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide)

    body_top = T.HEADER_H
    body_h = T.SLIDE_H - body_top
    block_h = Inches(3.4)
    block_top = body_top + (body_h - block_h) // 2

    accent_w = Inches(1.8)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        (T.SLIDE_W - accent_w) // 2, block_top,
        accent_w, Inches(0.06),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(T.ORANGE)
    accent.line.fill.background()

    _add_text(slide, Inches(0), block_top + Inches(0.35),
              T.SLIDE_W, Inches(1.6),
              "Thank you",
              font=T.FONT_HEAD, size=Pt(56), color=T.INK, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _add_text(slide, Inches(0), block_top + Inches(2.1),
              T.SLIDE_W, Inches(0.6),
              "Questions & discussion",
              font=T.FONT_BODY, size=Pt(22), color=T.MUTED, italic=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _add_text(slide, Inches(0), block_top + Inches(2.85),
              T.SLIDE_W, Inches(0.45),
              "speaker@kcl.ac.uk  ·  kclnlp.github.io",
              font=T.FONT_BODY, size=Pt(16), color=T.INK,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)


def build():
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H

    slide_title(prs)
    slide_title_centered(prs)
    slide_section(prs, number="01", title="Section heading")
    slide_content(prs, title="Slide title")
    slide_two_column(prs, title="Slide title")
    slide_thanks(prs)
    slide_thanks_centered(prs)

    prs.save(str(OUTPUT))
    print(f"wrote {OUTPUT}")


if __name__ == "__main__":
    build()
