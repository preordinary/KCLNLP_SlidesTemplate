"""
Build the "Paper Yellow" KCL NLP slide template.

Produces three files under the repo root:
    templates/paper-yellow/KCLNLP-PaperYellow.pptx      (blank, styled)
    templates/paper-yellow/KCLNLP-PaperYellow.potx      (same, as .potx)
    examples/paper-yellow/KCLNLP-PaperYellow-Example.pptx (populated sample)

Architecture
------------
We decouple two concerns so each works reliably:

1. **Slide Layouts** (in the template decks) — provide a proper PowerPoint
   "New Slide" UX: decoration shapes + typed placeholders where users fill
   in content. Placeholder inheritance works well in PowerPoint itself.

2. **Example deck slides** — drawn directly on blank-layout slides using
   static shapes + text boxes, to give us pixel-level control over how
   the preview renders (independent of whatever renderer is used).

Both paths share `decorate_N(container)` functions so the visual design
stays in one place.

Run:
    python scripts/build_paper_yellow.py
"""

from __future__ import annotations

import zipfile
from pathlib import Path

from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

# ---------- Palette ---------------------------------------------------------

HEX = {
    "paper":      "F6BD60",
    "cream":      "F7EDE2",
    "sage":       "84A59D",
    "terracotta": "F07167",
    "blush":      "FCD5CE",
    "ink":        "2B2B2B",
    "muted":      "6B6B6B",
    "white":      "FFFFFF",
}

# ---------- Fonts -----------------------------------------------------------

LATIN_MAJOR = "Aptos Display"
LATIN_MINOR = "Aptos"
EA_MAJOR    = "等线 Light"
EA_MINOR    = "等线"

# ---------- Geometry --------------------------------------------------------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------- Paths -----------------------------------------------------------

REPO = Path(__file__).resolve().parent.parent
LOGOS = REPO / "assets" / "logos"
LOGO_KCL    = LOGOS / "KCL.png"
LOGO_KCLNLP = LOGOS / "KCLNLP.png"
LOGO_ALAN   = LOGOS / "Alan.png"

OUT_TEMPLATE_DIR = REPO / "templates" / "paper-yellow"
OUT_EXAMPLE_DIR  = REPO / "examples" / "paper-yellow"
OUT_PPTX = OUT_TEMPLATE_DIR / "KCLNLP-PaperYellow.pptx"
OUT_POTX = OUT_TEMPLATE_DIR / "KCLNLP-PaperYellow.potx"
OUT_EXAMPLE = OUT_EXAMPLE_DIR / "KCLNLP-PaperYellow-Example.pptx"

WORDMARK = "King's College London NLP Group"

# ============================================================================
# XML namespaces
# ============================================================================

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def q(ns: str, tag: str) -> str:
    return f"{{{ns}}}{tag}"


# ============================================================================
# Container-agnostic helpers
# A "container" is anything with .shapes._spTree and .part — a SlideLayout
# or a Slide. We build XML and append to _spTree.
# ============================================================================

def _next_shape_id(container) -> int:
    used = []
    for c in container.shapes._spTree.iter():
        if "id" in c.attrib:
            try:
                used.append(int(c.attrib["id"]))
            except ValueError:
                pass
    return (max(used) + 1) if used else 10


def _append(container, xml: str) -> etree._Element:
    el = etree.fromstring(xml)
    nvCNvPr = el.find(".//" + q(P_NS, "cNvPr"))
    if nvCNvPr is not None and nvCNvPr.get("id") == "0":
        nvCNvPr.set("id", str(_next_shape_id(container)))
    container.shapes._spTree.append(el)
    return el


def _emu(v) -> int:
    return int(v)


def _xml_escape(s: str) -> str:
    return (s or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


# ============================================================================
# Layout-level background / name
# ============================================================================

def _strip(container) -> None:
    tree = container.shapes._spTree
    for child in list(tree):
        tag = etree.QName(child).localname
        if tag in {"sp", "pic", "grpSp", "cxnSp", "graphicFrame"}:
            tree.remove(child)


def _set_layout_name(layout, name: str) -> None:
    layout.element.cSld.set("name", name)


def _set_bg(container, hex_color: str) -> None:
    cSld = container.element.cSld
    for bg in cSld.findall(q(P_NS, "bg")):
        cSld.remove(bg)
    bg_xml = f'''<p:bg xmlns:p="{P_NS}" xmlns:a="{A_NS}">
  <p:bgPr>
    <a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>
    <a:effectLst/>
  </p:bgPr>
</p:bg>'''
    cSld.insert(0, etree.fromstring(bg_xml))


# ============================================================================
# Primitive shape builders
# ============================================================================

def add_rect(container, x, y, w, h, hex_fill):
    xml = f'''<p:sp xmlns:p="{P_NS}" xmlns:a="{A_NS}">
  <p:nvSpPr>
    <p:cNvPr id="0" name="Rect"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{_emu(x)}" y="{_emu(y)}"/>
      <a:ext cx="{_emu(w)}" cy="{_emu(h)}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="{hex_fill}"/></a:solidFill>
    <a:ln><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
</p:sp>'''
    return _append(container, xml)


def add_line(container, x1, y1, x2, y2, hex_color, weight_pt=1.5):
    x = min(x1, x2); y = min(y1, y2)
    cx = max(1, abs(x2 - x1)); cy = max(1, abs(y2 - y1))
    flipV = "1" if (y2 < y1) else "0"
    flipH = "1" if (x2 < x1) else "0"
    xml = f'''<p:cxnSp xmlns:p="{P_NS}" xmlns:a="{A_NS}">
  <p:nvCxnSpPr>
    <p:cNvPr id="0" name="Line"/>
    <p:cNvCxnSpPr/>
    <p:nvPr/>
  </p:nvCxnSpPr>
  <p:spPr>
    <a:xfrm flipH="{flipH}" flipV="{flipV}">
      <a:off x="{_emu(x)}" y="{_emu(y)}"/>
      <a:ext cx="{_emu(cx)}" cy="{_emu(cy)}"/>
    </a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:ln w="{int(weight_pt * 12700)}" cap="rnd">
      <a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>
      <a:round/>
    </a:ln>
  </p:spPr>
</p:cxnSp>'''
    return _append(container, xml)


def add_text(
    container, x, y, w, h, text, *,
    font=LATIN_MAJOR, size_pt=18, hex_color=HEX["ink"],
    bold=False, italic=False, align="l", anchor="t",
    line_spacing=None,
):
    """Add a text box. `text` may contain \\n for multiple paragraphs."""
    ea = EA_MAJOR if font == LATIN_MAJOR else EA_MINOR
    lines = _xml_escape(text).split("\n")
    line_sp_attr = ""
    if line_spacing is not None:
        line_sp_attr = f' <a:lnSpc><a:spcPct val="{int(line_spacing * 100000)}"/></a:lnSpc>'
    ps = []
    for line in lines:
        if line == "":
            ps.append('<a:p><a:pPr/><a:endParaRPr/></a:p>')
            continue
        ps.append(f'''<a:p>
  <a:pPr algn="{align}" marL="0" indent="0">{line_sp_attr}<a:buNone/></a:pPr>
  <a:r>
    <a:rPr lang="en-US" sz="{int(size_pt * 100)}" b="{1 if bold else 0}" i="{1 if italic else 0}" dirty="0">
      <a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>
      <a:latin typeface="{font}"/>
      <a:ea typeface="{ea}"/>
    </a:rPr>
    <a:t>{line}</a:t>
  </a:r>
</a:p>''')
    ps_xml = "\n".join(ps)
    xml = f'''<p:sp xmlns:p="{P_NS}" xmlns:a="{A_NS}">
  <p:nvSpPr>
    <p:cNvPr id="0" name="TextBox"/>
    <p:cNvSpPr txBox="1"/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{_emu(x)}" y="{_emu(y)}"/>
      <a:ext cx="{_emu(w)}" cy="{_emu(h)}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="square" anchor="{anchor}" lIns="0" rIns="0" tIns="0" bIns="0"/>
    <a:lstStyle/>
    {ps_xml}
  </p:txBody>
</p:sp>'''
    return _append(container, xml)


def add_bullets(
    container, x, y, w, h, items, *,
    font=LATIN_MINOR, size_pt=18, hex_color=HEX["ink"],
    bullet_color=None, bullet_char="•",
    line_spacing=1.25,
    anchor="t",
):
    """Add a text box with real bullet-style items."""
    ea = EA_MINOR if font == LATIN_MINOR else EA_MAJOR
    bullet_color = bullet_color or HEX["terracotta"]
    ps = []
    for item in items:
        safe = _xml_escape(item)
        ps.append(f'''<a:p>
  <a:pPr marL="342900" indent="-342900" algn="l">
    <a:lnSpc><a:spcPct val="{int(line_spacing * 100000)}"/></a:lnSpc>
    <a:buClr><a:srgbClr val="{bullet_color}"/></a:buClr>
    <a:buSzPct val="100000"/>
    <a:buFont typeface="{font}"/>
    <a:buChar char="{bullet_char}"/>
  </a:pPr>
  <a:r>
    <a:rPr lang="en-US" sz="{int(size_pt * 100)}" dirty="0">
      <a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>
      <a:latin typeface="{font}"/>
      <a:ea typeface="{ea}"/>
    </a:rPr>
    <a:t>{safe}</a:t>
  </a:r>
</a:p>''')
    ps_xml = "\n".join(ps)
    xml = f'''<p:sp xmlns:p="{P_NS}" xmlns:a="{A_NS}">
  <p:nvSpPr>
    <p:cNvPr id="0" name="Bullets"/>
    <p:cNvSpPr txBox="1"/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{_emu(x)}" y="{_emu(y)}"/>
      <a:ext cx="{_emu(w)}" cy="{_emu(h)}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="square" anchor="{anchor}" lIns="0" rIns="0" tIns="0" bIns="0"/>
    <a:lstStyle/>
    {ps_xml}
  </p:txBody>
</p:sp>'''
    return _append(container, xml)


def add_placeholder(
    container, *, idx: int, ph_type: str | None, name: str,
    x, y, w, h, prompt: str,
    font: str = LATIN_MAJOR, size_pt: int = 32,
    hex_color: str = HEX["ink"], bold: bool = False,
    align: str = "l", anchor: str = "t",
    no_bullet: bool = True,
):
    """Add a placeholder <p:sp> onto a layout. Strips default bullet list style."""
    ea = EA_MAJOR if font == LATIN_MAJOR else EA_MINOR
    type_attr = f' type="{ph_type}"' if ph_type else ""
    safe = _xml_escape(prompt)
    lines = safe.split("\n")
    ps = []
    for line in lines:
        ps.append(f'''<a:p>
  <a:pPr algn="{align}" marL="0" indent="0"><a:buNone/></a:pPr>
  <a:r>
    <a:rPr lang="en-US" sz="{int(size_pt * 100)}" b="{1 if bold else 0}" dirty="0">
      <a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>
      <a:latin typeface="{font}"/>
      <a:ea typeface="{ea}"/>
    </a:rPr>
    <a:t>{line}</a:t>
  </a:r>
</a:p>''')
    ps_xml = "\n".join(ps)
    # List style overrides inherited bullets (no_bullet=True) or leaves alone
    if no_bullet:
        lstStyle = f'''<a:lstStyle>
      <a:lvl1pPr marL="0" indent="0"><a:buNone/><a:defRPr/></a:lvl1pPr>
      <a:lvl2pPr marL="0" indent="0"><a:buNone/><a:defRPr/></a:lvl2pPr>
      <a:lvl3pPr marL="0" indent="0"><a:buNone/><a:defRPr/></a:lvl3pPr>
    </a:lstStyle>'''
    else:
        lstStyle = "<a:lstStyle/>"
    xml = f'''<p:sp xmlns:p="{P_NS}" xmlns:a="{A_NS}">
  <p:nvSpPr>
    <p:cNvPr id="0" name="{name}"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr><p:ph{type_attr} idx="{idx}"/></p:nvPr>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{_emu(x)}" y="{_emu(y)}"/>
      <a:ext cx="{_emu(w)}" cy="{_emu(h)}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="square" anchor="{anchor}" lIns="0" rIns="0" tIns="0" bIns="0"/>
    {lstStyle}
    {ps_xml}
  </p:txBody>
</p:sp>'''
    return _append(container, xml)


def add_slide_number_placeholder(container, x, y, w, h):
    xml = f'''<p:sp xmlns:p="{P_NS}" xmlns:a="{A_NS}">
  <p:nvSpPr>
    <p:cNvPr id="0" name="Slide Number"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr><p:ph type="sldNum" sz="quarter" idx="20"/></p:nvPr>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{_emu(x)}" y="{_emu(y)}"/>
      <a:ext cx="{_emu(w)}" cy="{_emu(h)}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr" lIns="0" rIns="0" tIns="0" bIns="0"/>
    <a:lstStyle><a:lvl1pPr algn="r" marL="0" indent="0"><a:buNone/></a:lvl1pPr></a:lstStyle>
    <a:p>
      <a:pPr algn="r" marL="0" indent="0"><a:buNone/></a:pPr>
      <a:fld id="{{B0C7D12E-1234-5678-90AB-CDEF01234567}}" type="slidenum">
        <a:rPr lang="en-US" sz="1000" dirty="0">
          <a:solidFill><a:srgbClr val="{HEX["sage"]}"/></a:solidFill>
          <a:latin typeface="{LATIN_MINOR}"/>
          <a:ea typeface="{EA_MINOR}"/>
        </a:rPr>
        <a:t>#</a:t>
      </a:fld>
    </a:p>
  </p:txBody>
</p:sp>'''
    return _append(container, xml)


def add_picture(container, path: Path, x, y, *, height=None, width=None):
    with Image.open(path) as im:
        px_w, px_h = im.size
    aspect = px_w / px_h
    if height is not None and width is None:
        w = int(height * aspect); h = int(height)
    elif width is not None and height is None:
        w = int(width); h = int(width / aspect)
    else:
        w = int(width); h = int(height)

    _, rId = container.part.get_or_add_image_part(str(path))
    name = path.stem
    xml = f'''<p:pic xmlns:p="{P_NS}" xmlns:a="{A_NS}" xmlns:r="{R_NS}">
  <p:nvPicPr>
    <p:cNvPr id="0" name="{name}"/>
    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{rId}"/>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="{_emu(x)}" y="{_emu(y)}"/>
      <a:ext cx="{_emu(w)}" cy="{_emu(h)}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>'''
    _append(container, xml)
    return w, h


# ============================================================================
# Shared decoration functions (work on layout OR slide)
# ============================================================================

def logo_stack_vertical(container, x, y, *, height, gap=Inches(0.2)):
    widths_heights = []
    for path in (LOGO_KCL, LOGO_KCLNLP, LOGO_ALAN):
        w, h = add_picture(container, path, x, y, height=height)
        widths_heights.append((w, h))
        y += h + gap


def logo_row_centered(container, y, *, height, gap=Inches(0.6)):
    widths = []
    for path in (LOGO_KCL, LOGO_KCLNLP, LOGO_ALAN):
        with Image.open(path) as im:
            aspect = im.width / im.height
        widths.append(int(height * aspect))
    total = sum(widths) + gap * (len(widths) - 1)
    x = (SLIDE_W - total) // 2
    for path, w in zip((LOGO_KCL, LOGO_KCLNLP, LOGO_ALAN), widths):
        add_picture(container, path, x, y, height=height)
        x += w + gap


def footer_strip(container, *, slide_number_placeholder=False, slide_number_text=None):
    """Small KCLNLP branding footer.

    - On layouts, pass slide_number_placeholder=True so PowerPoint renders the real #.
    - On example slides, pass slide_number_text="3" to show a static number.
    """
    y = Inches(7.05)
    w, h = add_picture(container, LOGO_KCLNLP, Inches(0.5), y, height=Inches(0.3))
    add_text(
        container, Inches(0.5) + w + Inches(0.15), y + Inches(0.03),
        Inches(7.0), Inches(0.3),
        WORDMARK,
        font=LATIN_MINOR, size_pt=10, hex_color=HEX["sage"],
    )
    if slide_number_placeholder:
        add_slide_number_placeholder(container, Inches(12.3), y, Inches(0.6), Inches(0.3))
    elif slide_number_text is not None:
        add_text(
            container, Inches(12.3), y, Inches(0.6), Inches(0.3),
            slide_number_text,
            font=LATIN_MINOR, size_pt=10, hex_color=HEX["sage"], align="r", anchor="ctr",
        )


# ============================================================================
# Per-layout DECORATION (static visual only — excludes placeholders and content)
# ============================================================================

def decorate_01_title(c):
    _set_bg(c, HEX["paper"])
    # Cream "paper" card on the right
    add_rect(c, Inches(8.2), Inches(0.8), Inches(4.5), Inches(5.9), HEX["cream"])
    # Terracotta accent square
    add_rect(c, Inches(8.0), Inches(0.6), Inches(0.5), Inches(0.5), HEX["terracotta"])
    # Sage accent square
    add_rect(c, Inches(12.6), Inches(6.5), Inches(0.3), Inches(0.3), HEX["sage"])
    # Logo stack on card
    logo_stack_vertical(c, Inches(8.6), Inches(1.15), height=Inches(0.8), gap=Inches(0.25))
    # Wordmark at bottom of card
    add_text(
        c, Inches(8.4), Inches(6.2), Inches(4.2), Inches(0.4),
        WORDMARK,
        font=LATIN_MINOR, size_pt=12, hex_color=HEX["sage"], bold=True, align="ctr",
    )


def decorate_02_toc(c):
    _set_bg(c, HEX["cream"])
    # Left paper-yellow accent bar
    add_rect(c, 0, 0, Inches(0.35), SLIDE_H, HEX["paper"])
    # Right column logo stack
    logo_stack_vertical(c, Inches(10.9), Inches(2.2), height=Inches(0.7), gap=Inches(0.3))


def decorate_03_section(c):
    _set_bg(c, HEX["cream"])
    # Paper-yellow left third
    add_rect(c, 0, 0, Inches(4.5), SLIDE_H, HEX["paper"])
    # Small sage square top-right
    add_rect(c, Inches(12.6), Inches(0.4), Inches(0.35), Inches(0.35), HEX["sage"])


def decorate_04_content_single(c):
    _set_bg(c, HEX["cream"])
    add_rect(c, 0, 0, Inches(0.35), SLIDE_H, HEX["paper"])


def decorate_05_content_two_col(c):
    _set_bg(c, HEX["cream"])
    add_rect(c, 0, 0, Inches(0.35), SLIDE_H, HEX["paper"])
    # Blush vertical separator between columns
    add_rect(c, Inches(6.665), Inches(1.9), Inches(0.04), Inches(4.8), HEX["blush"])


def decorate_06_closing(c):
    _set_bg(c, HEX["paper"])
    # Cream card
    add_rect(c, Inches(2.2), Inches(1.6), Inches(8.9), Inches(3.5), HEX["cream"])


# ============================================================================
# Layout builders (decoration + placeholders + (static wordmark/footer)
# ============================================================================

def build_layout_01(layout):
    _strip(layout)
    _set_layout_name(layout, "01 Title Slide")
    decorate_01_title(layout)

    # Title
    add_placeholder(
        layout, idx=0, ph_type="ctrTitle", name="Title",
        x=Inches(0.75), y=Inches(2.2), w=Inches(7.3), h=Inches(2.3),
        prompt="Presentation Title",
        font=LATIN_MAJOR, size_pt=54, hex_color=HEX["ink"], bold=True, anchor="b",
    )
    # Subtitle
    add_placeholder(
        layout, idx=1, ph_type="subTitle", name="Subtitle",
        x=Inches(0.75), y=Inches(4.55), w=Inches(7.3), h=Inches(0.9),
        prompt="A concise, descriptive subtitle",
        font=LATIN_MINOR, size_pt=22, hex_color=HEX["terracotta"], anchor="t",
    )
    # Author/date line as fixed text box
    add_text(
        layout, Inches(0.75), Inches(5.6), Inches(7.3), Inches(0.45),
        "Speaker Name · 2026",
        font=LATIN_MINOR, size_pt=16, hex_color=HEX["muted"],
    )


def build_layout_02(layout):
    _strip(layout)
    _set_layout_name(layout, "02 Table of Contents")
    decorate_02_toc(layout)

    add_placeholder(
        layout, idx=0, ph_type="title", name="Title",
        x=Inches(0.9), y=Inches(0.6), w=Inches(10), h=Inches(1.0),
        prompt="Contents",
        font=LATIN_MAJOR, size_pt=40, hex_color=HEX["ink"], bold=True,
    )
    add_line(layout, Inches(0.95), Inches(1.65), Inches(2.5), Inches(1.65), HEX["sage"], 2.0)
    add_placeholder(
        layout, idx=1, ph_type="body", name="Contents List",
        x=Inches(0.9), y=Inches(2.0), w=Inches(9.5), h=Inches(4.6),
        prompt="01   Introduction\n02   Related Work\n03   Method\n04   Experiments\n05   Conclusion",
        font=LATIN_MINOR, size_pt=24, hex_color=HEX["ink"],
    )
    footer_strip(layout, slide_number_placeholder=True)


def build_layout_03(layout):
    _strip(layout)
    _set_layout_name(layout, "03 Section Divider")
    decorate_03_section(layout)

    # Section number — static text on the yellow panel (not a placeholder)
    add_text(
        layout, Inches(0.7), Inches(2.3), Inches(3.5), Inches(2.2),
        "02",
        font=LATIN_MAJOR, size_pt=140, hex_color=HEX["terracotta"], bold=True,
        align="ctr", anchor="ctr",
    )
    # Section title
    add_placeholder(
        layout, idx=0, ph_type="title", name="Section Title",
        x=Inches(5.1), y=Inches(2.8), w=Inches(7.8), h=Inches(1.3),
        prompt="Section Title",
        font=LATIN_MAJOR, size_pt=44, hex_color=HEX["ink"], bold=True,
    )
    add_line(layout, Inches(5.1), Inches(4.15), Inches(6.5), Inches(4.15), HEX["sage"], 2.0)
    add_placeholder(
        layout, idx=1, ph_type="body", name="Section Lead-in",
        x=Inches(5.1), y=Inches(4.3), w=Inches(7.8), h=Inches(1.2),
        prompt="A one-line description of what this section is about.",
        font=LATIN_MINOR, size_pt=20, hex_color=HEX["muted"],
    )
    footer_strip(layout, slide_number_placeholder=True)


def build_layout_04(layout):
    _strip(layout)
    _set_layout_name(layout, "04 Content — Single Column")
    decorate_04_content_single(layout)

    add_placeholder(
        layout, idx=0, ph_type="title", name="Title",
        x=Inches(0.9), y=Inches(0.5), w=Inches(11.5), h=Inches(0.9),
        prompt="Slide Title",
        font=LATIN_MAJOR, size_pt=32, hex_color=HEX["ink"], bold=True,
    )
    add_line(layout, Inches(0.95), Inches(1.4), Inches(2.2), Inches(1.4), HEX["sage"], 2.0)
    add_placeholder(
        layout, idx=1, ph_type="body", name="Body",
        x=Inches(0.9), y=Inches(1.7), w=Inches(11.5), h=Inches(5.0),
        prompt="Click to add body text.",
        font=LATIN_MINOR, size_pt=20, hex_color=HEX["ink"],
    )
    footer_strip(layout, slide_number_placeholder=True)


def build_layout_05(layout):
    _strip(layout)
    _set_layout_name(layout, "05 Content — Two Column")
    decorate_05_content_two_col(layout)

    add_placeholder(
        layout, idx=0, ph_type="title", name="Title",
        x=Inches(0.9), y=Inches(0.5), w=Inches(11.5), h=Inches(0.9),
        prompt="Slide Title",
        font=LATIN_MAJOR, size_pt=32, hex_color=HEX["ink"], bold=True,
    )
    add_line(layout, Inches(0.95), Inches(1.4), Inches(2.2), Inches(1.4), HEX["sage"], 2.0)
    add_placeholder(
        layout, idx=1, ph_type="body", name="Left Column",
        x=Inches(0.9), y=Inches(1.9), w=Inches(5.6), h=Inches(4.9),
        prompt="Left column — context, motivation, or first side of comparison.",
        font=LATIN_MINOR, size_pt=18, hex_color=HEX["ink"],
    )
    add_placeholder(
        layout, idx=2, ph_type="body", name="Right Column",
        x=Inches(6.85), y=Inches(1.9), w=Inches(5.6), h=Inches(4.9),
        prompt="Right column — result, proposed approach, or second side.",
        font=LATIN_MINOR, size_pt=18, hex_color=HEX["ink"],
    )
    footer_strip(layout, slide_number_placeholder=True)


def build_layout_06(layout):
    _strip(layout)
    _set_layout_name(layout, "06 Closing / Thanks")
    decorate_06_closing(layout)

    add_placeholder(
        layout, idx=0, ph_type="ctrTitle", name="Thanks Title",
        x=Inches(2.4), y=Inches(1.9), w=Inches(8.5), h=Inches(1.8),
        prompt="Thank You",
        font=LATIN_MAJOR, size_pt=72, hex_color=HEX["terracotta"], bold=True,
        align="ctr", anchor="ctr",
    )
    add_placeholder(
        layout, idx=1, ph_type="subTitle", name="Contact",
        x=Inches(2.4), y=Inches(3.8), w=Inches(8.5), h=Inches(1.2),
        prompt="Questions? · name@kcl.ac.uk",
        font=LATIN_MINOR, size_pt=22, hex_color=HEX["ink"],
        align="ctr", anchor="t",
    )
    # Logos row centred below the card
    logo_row_centered(layout, Inches(5.55), height=Inches(0.85))
    # Wordmark at bottom (but above the preview-cut zone)
    add_text(
        layout, 0, Inches(6.85), SLIDE_W, Inches(0.4),
        WORDMARK,
        font=LATIN_MINOR, size_pt=12, hex_color=HEX["sage"], bold=True, align="ctr",
    )


# ============================================================================
# Example slides — draw everything directly on blank slides
# ============================================================================

def _blank_layout(prs):
    return prs.slide_layouts[6]  # Blank


def _add_blank_slide(prs):
    return prs.slides.add_slide(_blank_layout(prs))


def example_slide_01(prs):
    slide = _add_blank_slide(prs)
    decorate_01_title(slide)
    # Title
    add_text(
        slide, Inches(0.75), Inches(2.3), Inches(7.3), Inches(2.3),
        "Retrieval-Augmented Generation\nfor Scientific QA",
        font=LATIN_MAJOR, size_pt=40, hex_color=HEX["ink"], bold=True, anchor="b",
        line_spacing=1.1,
    )
    # Subtitle
    add_text(
        slide, Inches(0.75), Inches(4.75), Inches(7.3), Inches(0.7),
        "A pragmatic look at RAG for long-context research",
        font=LATIN_MINOR, size_pt=22, hex_color=HEX["terracotta"],
    )
    # Speaker / date
    add_text(
        slide, Inches(0.75), Inches(5.55), Inches(7.3), Inches(0.45),
        "Yizhen Yao · April 2026",
        font=LATIN_MINOR, size_pt=16, hex_color=HEX["muted"],
    )


def example_slide_02(prs):
    slide = _add_blank_slide(prs)
    decorate_02_toc(slide)
    add_text(
        slide, Inches(0.9), Inches(0.6), Inches(10), Inches(1.0),
        "Contents",
        font=LATIN_MAJOR, size_pt=44, hex_color=HEX["ink"], bold=True,
    )
    add_line(slide, Inches(0.95), Inches(1.65), Inches(2.5), Inches(1.65), HEX["sage"], 2.0)
    # Contents rows — numbered, no bullets
    rows = [
        ("01", "Motivation"),
        ("02", "Related Work"),
        ("03", "Method"),
        ("04", "Experiments"),
        ("05", "Takeaways"),
    ]
    y = Inches(2.2)
    for num, title in rows:
        add_text(
            slide, Inches(0.95), y, Inches(1.2), Inches(0.6),
            num,
            font=LATIN_MAJOR, size_pt=28, hex_color=HEX["terracotta"], bold=True,
        )
        add_text(
            slide, Inches(2.2), y + Inches(0.07), Inches(7.5), Inches(0.6),
            title,
            font=LATIN_MINOR, size_pt=24, hex_color=HEX["ink"],
        )
        y += Inches(0.85)
    footer_strip(slide, slide_number_text="2")


def example_slide_03(prs):
    slide = _add_blank_slide(prs)
    decorate_03_section(slide)
    # Large section number
    add_text(
        slide, Inches(0.4), Inches(2.3), Inches(4.0), Inches(2.4),
        "03",
        font=LATIN_MAJOR, size_pt=150, hex_color=HEX["terracotta"], bold=True,
        align="ctr", anchor="ctr",
    )
    # Section title right
    add_text(
        slide, Inches(5.1), Inches(2.9), Inches(7.8), Inches(1.3),
        "Method",
        font=LATIN_MAJOR, size_pt=48, hex_color=HEX["ink"], bold=True,
    )
    add_line(slide, Inches(5.1), Inches(4.15), Inches(6.5), Inches(4.15), HEX["sage"], 2.0)
    add_text(
        slide, Inches(5.1), Inches(4.3), Inches(7.8), Inches(1.2),
        "How we combine retrieval with long-context generation.",
        font=LATIN_MINOR, size_pt=20, hex_color=HEX["muted"],
    )
    footer_strip(slide, slide_number_text="3")


def example_slide_04(prs):
    slide = _add_blank_slide(prs)
    decorate_04_content_single(slide)
    add_text(
        slide, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.9),
        "Motivation",
        font=LATIN_MAJOR, size_pt=32, hex_color=HEX["ink"], bold=True,
    )
    add_line(slide, Inches(0.95), Inches(1.4), Inches(2.2), Inches(1.4), HEX["sage"], 2.0)
    add_bullets(
        slide, Inches(0.95), Inches(1.8), Inches(11.5), Inches(5.0),
        [
            "LLMs still hallucinate on specialised scientific questions.",
            "Retrieval narrows the input distribution at inference time.",
            "But naive RAG hurts on questions that need synthesis across documents.",
            "Goal: a lightweight recipe that keeps both grounding and synthesis.",
        ],
        font=LATIN_MINOR, size_pt=22, hex_color=HEX["ink"],
        bullet_color=HEX["terracotta"], line_spacing=1.4,
    )
    footer_strip(slide, slide_number_text="4")


def example_slide_05(prs):
    slide = _add_blank_slide(prs)
    decorate_05_content_two_col(slide)
    add_text(
        slide, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.9),
        "Baseline vs. Ours",
        font=LATIN_MAJOR, size_pt=32, hex_color=HEX["ink"], bold=True,
    )
    add_line(slide, Inches(0.95), Inches(1.4), Inches(2.2), Inches(1.4), HEX["sage"], 2.0)

    # Left column heading
    add_text(
        slide, Inches(0.95), Inches(1.9), Inches(5.6), Inches(0.55),
        "Baseline RAG",
        font=LATIN_MAJOR, size_pt=22, hex_color=HEX["sage"], bold=True,
    )
    add_bullets(
        slide, Inches(0.95), Inches(2.5), Inches(5.6), Inches(4.3),
        [
            "Top-k dense retrieval",
            "Stuffed into the prompt",
            "No reranking or clustering",
            "Answers tend to be extractive",
        ],
        font=LATIN_MINOR, size_pt=18, hex_color=HEX["ink"],
        bullet_color=HEX["sage"], line_spacing=1.4,
    )

    # Right column heading
    add_text(
        slide, Inches(6.9), Inches(1.9), Inches(5.6), Inches(0.55),
        "Hierarchical RAG (ours)",
        font=LATIN_MAJOR, size_pt=22, hex_color=HEX["terracotta"], bold=True,
    )
    add_bullets(
        slide, Inches(6.9), Inches(2.5), Inches(5.6), Inches(4.3),
        [
            "Cluster-aware retrieval",
            "Summary sketch + evidence",
            "Tool-use for citations",
            "+6.8 F1, −32% hallucinations",
        ],
        font=LATIN_MINOR, size_pt=18, hex_color=HEX["ink"],
        bullet_color=HEX["terracotta"], line_spacing=1.4,
    )
    footer_strip(slide, slide_number_text="5")


def example_slide_06(prs):
    slide = _add_blank_slide(prs)
    decorate_06_closing(slide)
    add_text(
        slide, Inches(2.4), Inches(2.1), Inches(8.5), Inches(1.8),
        "Thank You",
        font=LATIN_MAJOR, size_pt=72, hex_color=HEX["terracotta"], bold=True,
        align="ctr", anchor="ctr",
    )
    add_text(
        slide, Inches(2.4), Inches(3.9), Inches(8.5), Inches(0.8),
        "Questions? · yizhen.yao@kcl.ac.uk",
        font=LATIN_MINOR, size_pt=22, hex_color=HEX["ink"], align="ctr",
    )
    logo_row_centered(slide, Inches(5.55), height=Inches(0.85))
    add_text(
        slide, 0, Inches(6.85), SLIDE_W, Inches(0.4),
        WORDMARK,
        font=LATIN_MINOR, size_pt=12, hex_color=HEX["sage"], bold=True, align="ctr",
    )


EXAMPLE_BUILDERS = [
    example_slide_01, example_slide_02, example_slide_03,
    example_slide_04, example_slide_05, example_slide_06,
]


# ============================================================================
# Theme XML patching
# ============================================================================

def _theme_color_scheme_xml() -> str:
    return f'''<a:clrScheme name="Paper Yellow" xmlns:a="{A_NS}">
    <a:dk1><a:srgbClr val="{HEX["ink"]}"/></a:dk1>
    <a:lt1><a:srgbClr val="{HEX["white"]}"/></a:lt1>
    <a:dk2><a:srgbClr val="2B2B2B"/></a:dk2>
    <a:lt2><a:srgbClr val="{HEX["cream"]}"/></a:lt2>
    <a:accent1><a:srgbClr val="{HEX["paper"]}"/></a:accent1>
    <a:accent2><a:srgbClr val="{HEX["terracotta"]}"/></a:accent2>
    <a:accent3><a:srgbClr val="{HEX["sage"]}"/></a:accent3>
    <a:accent4><a:srgbClr val="{HEX["blush"]}"/></a:accent4>
    <a:accent5><a:srgbClr val="{HEX["cream"]}"/></a:accent5>
    <a:accent6><a:srgbClr val="AD8A56"/></a:accent6>
    <a:hlink><a:srgbClr val="{HEX["terracotta"]}"/></a:hlink>
    <a:folHlink><a:srgbClr val="{HEX["sage"]}"/></a:folHlink>
  </a:clrScheme>'''


def _theme_font_scheme_xml() -> str:
    return f'''<a:fontScheme name="KCLNLP Paper" xmlns:a="{A_NS}">
    <a:majorFont>
      <a:latin typeface="{LATIN_MAJOR}"/>
      <a:ea typeface=""/>
      <a:cs typeface=""/>
      <a:font script="Hans" typeface="{EA_MAJOR}"/>
      <a:font script="Hant" typeface="新細明體"/>
      <a:font script="Jpan" typeface="游ゴシック Light"/>
      <a:font script="Hang" typeface="맑은 고딕"/>
    </a:majorFont>
    <a:minorFont>
      <a:latin typeface="{LATIN_MINOR}"/>
      <a:ea typeface=""/>
      <a:cs typeface=""/>
      <a:font script="Hans" typeface="{EA_MINOR}"/>
      <a:font script="Hant" typeface="新細明體"/>
      <a:font script="Jpan" typeface="游ゴシック"/>
      <a:font script="Hang" typeface="맑은 고딕"/>
    </a:minorFont>
  </a:fontScheme>'''


def patch_theme_in_pptx(pptx_path: Path) -> None:
    tmp = pptx_path.with_suffix(".patching.pptx")
    with zipfile.ZipFile(pptx_path, "r") as zin, zipfile.ZipFile(
        tmp, "w", zipfile.ZIP_DEFLATED
    ) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "ppt/theme/theme1.xml":
                data = _patched_theme(data)
            zout.writestr(item, data)
    tmp.replace(pptx_path)


def _patched_theme(xml_bytes: bytes) -> bytes:
    tree = etree.fromstring(xml_bytes)
    theme_elts = tree.find(q(A_NS, "themeElements"))
    for old in theme_elts.findall(q(A_NS, "clrScheme")):
        theme_elts.remove(old)
    theme_elts.insert(0, etree.fromstring(_theme_color_scheme_xml()))
    for old in theme_elts.findall(q(A_NS, "fontScheme")):
        theme_elts.remove(old)
    theme_elts.insert(1, etree.fromstring(_theme_font_scheme_xml()))
    return etree.tostring(tree, xml_declaration=True, encoding="UTF-8", standalone=True)


# ============================================================================
# .potx variant
# ============================================================================

CT_PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
CT_POTX = "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"


def save_as_potx(pptx_path: Path, potx_path: Path) -> None:
    potx_path.parent.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(pptx_path, "r") as zin, zipfile.ZipFile(
        potx_path, "w", zipfile.ZIP_DEFLATED
    ) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                data = data.replace(CT_PPTX.encode(), CT_POTX.encode())
            zout.writestr(item, data)


# ============================================================================
# Build orchestration
# ============================================================================

# Default python-pptx layouts: 0=Title, 1=Title+Content, 2=Section Header,
# 3=Two Content, 4=Comparison, 5=Title Only, 6=Blank, 7+=Content w/Caption etc.
LAYOUT_PLAN = [
    (0, build_layout_01),
    (5, build_layout_02),
    (2, build_layout_03),
    (1, build_layout_04),
    (3, build_layout_05),
    (6, build_layout_06),  # NB: we also use this as "Blank" for example slides
]


def build_presentation(*, populate_example: bool = False) -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Customise the 6 layouts we care about (except we skip customising idx=6
    # when populate_example=True because example slides need a true Blank layout)
    used = set()
    for idx, builder in LAYOUT_PLAN:
        if populate_example and idx == 6:
            # Leave idx=6 as raw Blank; the example deck draws everything on it.
            # But still rename so users looking at the example pptx see a sensible layout list.
            layout = prs.slide_layouts[idx]
            _strip(layout)
            _set_layout_name(layout, "Blank")
            used.add(idx)
            continue
        layout = prs.slide_layouts[idx]
        builder(layout)
        used.add(idx)

    # Rename unused layouts so they sort last
    aux_n = 1
    for i, layout in enumerate(prs.slide_layouts):
        if i in used:
            continue
        _strip(layout)
        _set_layout_name(layout, f"_Aux {aux_n}")
        _set_bg(layout, HEX["cream"])
        aux_n += 1

    if populate_example:
        for builder in EXAMPLE_BUILDERS:
            builder(prs)

    return prs


# ============================================================================
# Main
# ============================================================================

def main() -> None:
    OUT_TEMPLATE_DIR.mkdir(parents=True, exist_ok=True)
    OUT_EXAMPLE_DIR.mkdir(parents=True, exist_ok=True)

    prs = build_presentation(populate_example=False)
    prs.save(OUT_PPTX)
    patch_theme_in_pptx(OUT_PPTX)
    print(f"Wrote {OUT_PPTX.relative_to(REPO)}")

    save_as_potx(OUT_PPTX, OUT_POTX)
    print(f"Wrote {OUT_POTX.relative_to(REPO)}")

    prs2 = build_presentation(populate_example=True)
    prs2.save(OUT_EXAMPLE)
    patch_theme_in_pptx(OUT_EXAMPLE)
    print(f"Wrote {OUT_EXAMPLE.relative_to(REPO)}")


if __name__ == "__main__":
    main()
